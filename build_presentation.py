"""
Genera INFORME_PRESENTACION.pptx a partir de los contenidos de INFORME.md
y las imágenes de informe_assets/.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

ASSETS = "informe_assets"

# Ruta al logo institucional (UBA / FIUBA). Si no existe, se omite sin error.
LOGO_PATH = "logo-fiuba.png"

# ---------------------------------------------------------------- paleta (tema claro y moderno)
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)
ALT_BG = RGBColor(0xE3, 0xEC, 0xF7)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0xA8, 0x3C)
BORDER = RGBColor(0xC9, 0xD4, 0xE3)
GOOD_BG = RGBColor(0xD6, 0xEA, 0xD2)
HEADER_BG = RGBColor(0xEC, 0xF2, 0xFB)
PAGE_BG = RGBColor(0xFC, 0xFD, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text):
    """Agrega notas del orador a la diapositiva."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.word_wrap = True
    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        line = line.strip()
        if i == 0:
            p = notes_tf.paragraphs[0]
        else:
            p = notes_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)


def add_logo(slide, left, top, height):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


def add_footer(slide, page_num):
    # línea fina separadora, estilo moderno (sin barra sólida pesada)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.32), SW, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), SH - Inches(0.32), SW - Inches(1.6), Inches(0.3))
    tf = box.text_frame
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Generación Eléctrica Mensual de Argentina (2015-2023)  ·  Análisis de Series de Tiempo — 02MIA2026"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

    pg = slide.shapes.add_textbox(SW - Inches(1.0), SH - Inches(0.32), Inches(0.75), Inches(0.3))
    tf2 = pg.text_frame
    tf2.margin_top = Inches(0.03)
    tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title, kicker=None):
    # banda superior clara, look moderno con borde inferior fino y filete de acento
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Pt(1.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), Inches(1.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(2.4), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    if kicker:
        p0 = tf.paragraphs[0]
        p0.text = kicker
        p0.font.size = Pt(14)
        p0.font.color.rgb = TEAL
        p0.font.bold = True
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # logo institucional arriba a la derecha, sobre la banda clara
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, SW - Inches(1.9), Inches(0.18), height=Inches(0.8))


def add_bullets(slide, items, left, top, width, height, font_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if text == "":
            p.text = ""
            p.space_after = Pt(4)
            continue
        p.text = ("•  " if level == 0 else "‒  ") + text
        p.level = level
        p.font.size = Pt(font_size if level == 0 else font_size - 2)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33) if level == 0 else GRAY
        p.space_after = Pt(8)
    return box


def add_picture_fit(slide, img_path, left, top, max_w, max_h):
    im = Image.open(img_path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(img_path, x, y, width=w, height=h)


def set_cell_border(cell, color=BORDER, width_pt=0.75):
    """Aplica un borde fino a las 4 aristas de una celda de tabla."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    width_emu = Pt(width_pt)
    hexcolor = "%02X%02X%02X" % (color[0], color[1], color[2])
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
        ln = tcPr.makeelement(qn(tag), {"w": str(int(width_emu)), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        srgb = fill.makeelement(qn("a:srgbClr"), {"val": hexcolor})
        fill.append(srgb)
        ln.append(fill)
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
        tcPr.append(ln)


def style_table(table, header_color=NAVY, font_size=14, header_font_size=15, col_widths=None,
                 row_height=None, header_height=None):
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    if row_height:
        for r in range(1, len(table.rows)):
            table.rows[r].height = row_height
    if header_height:
        table.rows[0].height = header_height
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 or r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(header_font_size if r == 0 else font_size)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = RGBColor(0x2B, 0x2B, 0x2B)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_BG if r % 2 == 0 else WHITE
            set_cell_border(cell)


def fill_table_text(table, data):
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)


def highlight_row(table, row_idx, color=GOOD_BG):
    for c in range(len(table.columns)):
        cell = table.cell(row_idx, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True


# ============================================================== SLIDE 1 — TÍTULO
slide = add_slide()
# fondo general claro
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = PAGE_BG
bg.line.fill.background()

# banda lateral izquierda en navy, con elementos decorativos modernos
side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.3), SH)
side.fill.solid()
side.fill.fore_color.rgb = NAVY
side.line.fill.background()

circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.3), Inches(-1.3), Inches(3.2), Inches(3.2))
circle1.fill.solid()
circle1.fill.fore_color.rgb = TEAL
circle1.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.6), Inches(5.6), Inches(2.6), Inches(2.6))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT
circle2.line.fill.background()

# placa blanca con el logo institucional, sobre la banda navy
if os.path.exists(LOGO_PATH):
    plate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(2.6), Inches(1.05))
    plate.fill.solid()
    plate.fill.fore_color.rgb = WHITE
    plate.line.fill.background()
    plate.shadow.inherit = False
    add_picture_fit(slide, LOGO_PATH, Inches(0.7), Inches(0.65), Inches(2.3), Inches(0.85))

tag = slide.shapes.add_textbox(Inches(0.55), Inches(5.55), Inches(3.4), Inches(1.6))
tftag = tag.text_frame
tftag.word_wrap = True
ptag = tftag.paragraphs[0]
ptag.text = "Especialización en\nInteligencia Artificial"
ptag.font.size = Pt(16)
ptag.font.bold = True
ptag.font.color.rgb = WHITE
ptag2 = tftag.add_paragraph()
ptag2.text = "Facultad de Ingeniería — UBA"
ptag2.font.size = Pt(14)
ptag2.font.color.rgb = RGBColor(0xCF, 0xD9, 0xEC)
ptag2.space_before = Pt(6)

# título principal, lado derecho sobre fondo claro
tb = slide.shapes.add_textbox(Inches(4.85), Inches(1.0), Inches(8.0), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Análisis de Series de Tiempo Aplicado a la"
p.font.size = Pt(24)
p.font.color.rgb = TEAL
p2 = tf.add_paragraph()
p2.text = "Generación Eléctrica Mensual de Argentina"
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = NAVY
p3 = tf.add_paragraph()
p3.text = "(2015 – 2023)"
p3.font.size = Pt(22)
p3.font.color.rgb = ACCENT
p3.font.bold = True
p3.space_before = Pt(6)

# tarjeta con los datos de la materia, docente y alumnos
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(4.35), Inches(7.9), Inches(2.45))
card.fill.solid()
card.fill.fore_color.rgb = WHITE
card.line.color.rgb = BORDER
card.line.width = Pt(1)
card.shadow.inherit = False

tb2 = card.text_frame
tb2.word_wrap = True
tb2.margin_left = Inches(0.35)
tb2.margin_right = Inches(0.35)
tb2.margin_top = Inches(0.25)

info_lines = [
    ("Curso:  ", "Análisis de Series de Tiempo (02MIA2026)"),
    ("Programa:  ", "Especialización en Inteligencia Artificial — UBA, FIUBA"),
    ("Docente:  ", "Camilo Argoty"),
    ("Alumnos:  ", "Gustavo Rivas · Carlos Rivas · Fermín Rodríguez"),
]
first = True
for label, value in info_lines:
    if first:
        p = tb2.paragraphs[0]
        first = False
    else:
        p = tb2.add_paragraph()
    p.font.size = Pt(17)
    p.space_after = Pt(8)
    r1 = p.add_run()
    r1.text = label
    r1.font.bold = True
    r1.font.color.rgb = TEAL
    r2 = p.add_run()
    r2.text = value
    r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

set_notes(slide, """
Buenas tardes. Vamos a presentar nuestro trabajo final de Análisis de Series de Tiempo,
donde aplicamos distintas técnicas de modelado y pronóstico sobre la generación eléctrica
mensual de Argentina entre 2015 y 2023.
Somos Gustavo Rivas, Carlos Rivas y Fermín Rodríguez, y el trabajo fue dirigido por el
profesor Camilo Argoty en el marco de la Especialización en Inteligencia Artificial de la UBA.
A lo largo de la presentación vamos a mostrar cómo analizamos los datos, qué modelos probamos
y cuál resultó más preciso para predecir la generación eléctrica.
""")

# ============================================================== SLIDE 2 — AGENDA
slide = add_slide()
add_title_bar(slide, "Agenda", "Trabajo Final")
agenda_items_l = [
    "1. De qué se trata el trabajo y qué queremos responder",
    "2. Los datos: de dónde vienen y cómo son",
    "3. Qué patrones encontramos (tendencia, estacionalidad)",
    "4. Los modelos que probamos: ARIMA, SARIMA, Holt-Winters y Prophet",
]
agenda_items_r = [
    "5. Cómo comparamos los modelos y cuál ganó",
    "6. Validación: ¿el modelo ganador está bien armado?",
    "7. Un extra: medir la volatilidad con GARCH",
    "8. Pronóstico para 2024 y conclusiones",
]
add_bullets(slide, agenda_items_l, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.5), font_size=22)
add_bullets(slide, agenda_items_r, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.5), font_size=22)
add_footer(slide, 2)
set_notes(slide, """
Así está organizada la presentación: primero contamos qué pregunta nos motivó y de dónde
sacamos los datos. Después mostramos qué patrones encontramos en la serie -tendencia y
estacionalidad-, y presentamos los cuatro modelos que construimos.
Luego viene la parte central: comparamos el desempeño de los modelos, validamos que el
ganador esté bien especificado, y agregamos dos extensiones -selección automática de
órdenes y un modelo de volatilidad GARCH-. Cerramos con el pronóstico para 2024 y las
conclusiones generales.
""")

# ============================================================== SLIDE 3 — PREGUNTA DE INVESTIGACIÓN
slide = add_slide()
add_title_bar(slide, "¿Qué queremos responder?", "1. Pregunta de Investigación")
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.7), Inches(1.55))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xEA, 0xF1, 0xFA)
box.line.color.rgb = TEAL
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_right = Inches(0.3)
p = tf.paragraphs[0]
p.text = ("¿Qué modelo predice mejor la generación eléctrica mensual de Argentina: "
          "ARIMA, SARIMA, Holt-Winters o Prophet? ¿Y qué tan importante es la fuerte "
          "estacionalidad anual -con picos en verano e invierno- para el resultado?")
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.5))
sub.text_frame.paragraphs[0].text = "En otras palabras, queremos lograr tres cosas:"
sub.text_frame.paragraphs[0].font.size = Pt(20)
sub.text_frame.paragraphs[0].font.bold = True
sub.text_frame.paragraphs[0].font.color.rgb = TEAL

obj_items = [
    "Entender cómo se mueve la generación eléctrica: ¿tiene tendencia? ¿se repite un patrón cada año?",
    "Probar cuatro modelos distintos y comparar qué tan bien predicen, apoyándonos en pruebas estadísticas formales para justificar cada decisión (no \"a ojo\").",
    "Elegir el mejor modelo y usarlo para proyectar la generación eléctrica de 2024, con un rango de confianza.",
]
add_bullets(slide, obj_items, Inches(0.9), Inches(3.8), Inches(11.5), Inches(3.0), font_size=18)
add_footer(slide, 3)
set_notes(slide, """
La pregunta central del trabajo es bastante simple de entender: queremos saber cuál de
estos cuatro modelos predice mejor la generación eléctrica mes a mes en Argentina, y
sobre todo, cuánto importa la estacionalidad -es decir, esos picos que se repiten todos
los veranos e inviernos-.
Para responderla nos propusimos tres cosas: primero entender los patrones de la serie
(tendencia y estacionalidad), después poner a competir los cuatro modelos con pruebas
estadísticas que justifiquen cada decisión, y finalmente usar el ganador para proyectar
2024 con un margen de confianza.
""")

# ============================================================== SLIDE 4 — MARCO TEÓRICO
slide = add_slide()
add_title_bar(slide, "Dos formas distintas de \"mirar\" una serie", "1. Marco Teórico")
left_items = [
    "Familia 1 — ARIMA / SARIMA:",
    ("Necesitan que la serie sea \"estable\" (estacionaria): sin tendencias raras ni cambios de comportamiento en el tiempo", 1),
    ("Si no lo es, se la transforma (diferenciando) hasta que lo sea", 1),
    ("Aprenden la dependencia entre un mes y los anteriores", 1),
    "",
    "Familia 2 — Holt-Winters / Prophet:",
    ("No exigen esa estabilidad previa", 1),
    ("Separan directamente la serie en \"nivel\", \"tendencia\" y \"estacionalidad\", y los van actualizando mes a mes", 1),
]
add_bullets(slide, left_items, Inches(0.8), Inches(1.5), Inches(7.0), Inches(4.6), font_size=18)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.5), Inches(4.5), Inches(4.7))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xE2)
box.line.color.rgb = ACCENT
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.25)
tf.margin_right = Inches(0.25)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "Nuestra hipótesis"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = ("Como la generación eléctrica argentina tiene un patrón estacional muy marcado, "
           "esperamos que los modelos que lo capturan explícitamente (SARIMA, Holt-Winters, "
           "Prophet) le ganen por bastante a un modelo que no lo tiene en cuenta (ARIMA).")
p2.font.size = Pt(16)
p2.space_before = Pt(10)
p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
p3 = tf.add_paragraph()
p3.text = ("Además, notamos que algunos meses son más \"volátiles\" que otros, así que "
           "sumamos un modelo extra (GARCH) para medir esa variabilidad.")
p3.font.size = Pt(16)
p3.space_before = Pt(10)
p3.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
add_footer(slide, 4)
set_notes(slide, """
Para entender por qué elegimos estos cuatro modelos, conviene pensarlos en dos grandes
familias. La primera son ARIMA y SARIMA, que vienen de la estadística clásica: necesitan
que la serie sea "estable" en el tiempo, y si no lo es, se la transforma diferenciándola.
La segunda familia es Holt-Winters y Prophet, que no exigen esa estabilidad previa: separan
directamente la serie en nivel, tendencia y estacionalidad, y van actualizando esos
componentes a medida que llegan nuevos datos.
Nuestra hipótesis de partida es simple: como la demanda eléctrica argentina tiene un
patrón estacional muy marcado -veranos e inviernos-, esperamos que los modelos que lo
capturan directamente le ganen claramente a ARIMA, que no lo modela. Y como notamos que
hay meses más volátiles que otros, agregamos GARCH como un análisis extra para medir esa
variabilidad.
""")

# ============================================================== SLIDE 5 — ORIGEN Y ATRIBUTOS DE LOS DATOS
slide = add_slide()
add_title_bar(slide, "De dónde vienen nuestros datos", "2. Descripción de los Datos")
items = [
    "Dataset público de Kaggle: \"Electricity Production Dataset\" (autor: sazidthe1)",
    ("Trae datos mensuales de generación eléctrica de 48 países, entre 2010 y 2023", 1),
    ("Nos quedamos solo con Argentina y con la generación neta total (\"Net Electricity Production\")", 1),
    ("Resultado: una serie de 108 meses, de enero 2015 a diciembre 2023", 1),
    "",
    "Después de filtrar, la serie con la que trabajamos es muy simple:",
]
add_bullets(slide, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.6), font_size=19)

tbl_shape = slide.shapes.add_table(3, 2, Inches(1.8), Inches(4.25), Inches(9.7), Inches(1.5))
table = tbl_shape.table
fill_table_text(table, [
    ["Columna", "Qué significa"],
    ["Fecha (mes)", "Un dato por mes, de enero de 2015 a diciembre de 2023"],
    ["Generación (GWh)", "Cuánta electricidad se generó en Argentina ese mes, en gigavatios-hora"],
])
style_table(table, font_size=17, header_font_size=18, col_widths=[Inches(2.7), Inches(7.0)],
            header_height=Inches(0.6), row_height=Inches(0.55))
add_footer(slide, 5)
set_notes(slide, """
Los datos vienen de un dataset público de Kaggle que tiene la generación eléctrica
mensual de 48 países. De ahí nos quedamos solo con Argentina, y dentro de los distintos
indicadores que trae, usamos la generación neta total -para no mezclar ni duplicar
valores por tipo de fuente de energía-.
El resultado es una serie muy simple de leer: un valor de generación eléctrica, en
gigavatios-hora, por cada uno de los 108 meses entre enero de 2015 y diciembre de 2023.
Esa es la serie univariada sobre la que trabajamos todo el análisis.
""")

# ============================================================== SLIDE 6 — CALIDAD DE DATOS
slide = add_slide()
add_title_bar(slide, "¿Los datos están \"limpios\"?", "2.3 Calidad de los Datos")
items = [
    "No hay valores faltantes",
    "No falta ningún mes: la serie está completa y es continua",
    "No detectamos valores atípicos (outliers)",
    "",
    "Conclusión: no hicieron falta correcciones, la serie está lista para analizarse tal cual.",
]
add_bullets(slide, items, Inches(0.8), Inches(1.5), Inches(5.6), Inches(3.2), font_size=20)

tbl_shape = slide.shapes.add_table(8, 2, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.6))
table = tbl_shape.table
fill_table_text(table, [
    ["Estadístico", "Valor (GWh)"],
    ["Cantidad de meses", "108"],
    ["Promedio", "33.811"],
    ["Desvío estándar", "3.058"],
    ["Mínimo", "25.656"],
    ["Percentil 25", "31.450"],
    ["Mediana", "33.508"],
    ["Percentil 75 / Máximo", "35.881 / 40.883"],
])
style_table(table, font_size=16, header_font_size=17, col_widths=[Inches(3.0), Inches(2.6)],
            header_height=Inches(0.55), row_height=Inches(0.5))

note = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(5.6), Inches(1.8))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("La variación entre meses representa solo ~9% del promedio: "
          "es decir, la serie no \"salta\" de forma errática, gran parte de ese "
          "movimiento se explica por la época del año (estacionalidad).")
p.font.size = Pt(17)
p.font.italic = True
p.font.color.rgb = GRAY
add_footer(slide, 6)
set_notes(slide, """
Antes de modelar, lo primero es chequear que los datos estén bien. Y en este caso la
buena noticia es que la serie viene muy limpia: no hay meses faltantes, no hay valores
nulos, y al aplicar el método estándar de detección de outliers no encontramos ninguno.
Mirando las estadísticas descriptivas, la generación promedio es de unos 33.800 GWh por
mes, con una variación de alrededor del 9% respecto a ese promedio. Esa variación
relativamente moderada es una primera pista de que el comportamiento de la serie está
dominado por un patrón que se repite -la estacionalidad- y no por saltos erráticos.
""")

# ============================================================== SLIDE 7 — SERIE ORIGINAL
slide = add_slide()
add_title_bar(slide, "Así se ve la serie completa", "2.4 Primer Vistazo a los Datos")
add_picture_fit(slide, f"{ASSETS}/13_serie_original.png", Inches(0.6), Inches(1.35), Inches(8.3), Inches(4.3))
items = [
    "Nivel promedio: ~33.800 GWh por mes",
    "No hay una tendencia clara de crecimiento o caída en los 9 años",
    "Pero sí hay una baja marcada en los últimos años: -11,6% en 2022 y -19,3% en 2023",
    "Se repite todos los años el mismo patrón: sube en verano e invierno",
    "Se ven caídas puntuales en momentos puntuales (ej. pandemia 2020)",
]
add_bullets(slide, items, Inches(9.1), Inches(1.6), Inches(4.0), Inches(4.8), font_size=16)
add_footer(slide, 7)
set_notes(slide, """
Este es el primer gráfico que armamos, y ya cuenta gran parte de la historia. Se ve que
la serie oscila siempre alrededor de un nivel promedio de unos 33.800 GWh por mes, sin
una tendencia clara de crecimiento o caída a lo largo de los nueve años -de hecho hicimos
una prueba estadística (correlación de Spearman) que confirma que no hay una tendencia
significativa-.
Lo que sí se nota a simple vista son dos cosas: primero, un patrón que se repite todos
los años, con subidas en verano e invierno; y segundo, una caída bastante marcada en los
últimos dos años -2022 y 2023-, que más adelante vamos a ver que le complica la vida a
todos los modelos al momento de evaluar.
""")

# ============================================================== SLIDE 8 — TASAS DE CRECIMIENTO
slide = add_slide()
add_title_bar(slide, "¿Cuánto creció (o cayó) cada año?", "2.5 Tasas de Crecimiento")
add_picture_fit(slide, f"{ASSETS}/16_tasas_crecimiento.png", Inches(0.6), Inches(1.35), Inches(7.6), Inches(4.5))

tbl_shape = slide.shapes.add_table(10, 2, Inches(8.6), Inches(1.35), Inches(4.1), Inches(4.6))
table = tbl_shape.table
fill_table_text(table, [
    ["Año", "Variación"],
    ["2015", "+0,07 %"],
    ["2016", "-2,93 %"],
    ["2017", "-5,38 %"],
    ["2018", "-14,06 %"],
    ["2019", "+0,87 %"],
    ["2020", "+2,10 %"],
    ["2021", "-0,46 %"],
    ["2022", "-11,64 %"],
    ["2023", "-19,26 %"],
])
style_table(table, font_size=14, header_font_size=15, col_widths=[Inches(1.9), Inches(2.2)],
            header_height=Inches(0.45), row_height=Inches(0.42))
highlight_row(table, 8, RGBColor(0xFC, 0xE4, 0xCB))
highlight_row(table, 9, RGBColor(0xFC, 0xE4, 0xCB))
add_footer(slide, 8)
set_notes(slide, """
Esta tabla muestra cuánto varió la generación eléctrica de un año a otro. Se ve que en
la mayoría de los años los cambios son moderados, para arriba o para abajo. Pero en los
dos últimos años -2022 y 2023, resaltados en naranja- la caída es mucho más fuerte:
-11,6% y -19,3%.
Esto es importante porque justamente esos dos años -sobre todo 2023- son los que vamos a
usar para evaluar qué tan bien predicen los modelos. Como ningún modelo "sabe" de
antemano que viene esa caída tan fuerte, todos van a tener cierto error en ese tramo.
Además, esta variación mes a mes no es siempre igual: hay períodos más tranquilos y
otros más movidos, algo que retomamos más adelante con el modelo GARCH.
""")

# ============================================================== SLIDE 9 — DESCOMPOSICIÓN
slide = add_slide()
add_title_bar(slide, "Separando la serie en sus \"piezas\"", "3.1 Tendencia, Estacionalidad y Residuo")
add_picture_fit(slide, f"{ASSETS}/19_descomposicion_aditiva.png", Inches(0.4), Inches(1.3), Inches(6.2), Inches(4.5))
add_picture_fit(slide, f"{ASSETS}/21_descomposicion_multiplicativa.png", Inches(6.7), Inches(1.3), Inches(6.2), Inches(4.5))
note = slide.shapes.add_textbox(Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.2))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Toda serie se puede pensar como Tendencia + Estacionalidad + Residuo (esquema "
          "\"aditivo\", izquierda) o como Tendencia × Estacionalidad × Residuo (\"multiplicativo\", derecha).")
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p2 = tf.add_paragraph()
p2.text = ("En nuestro caso, ambos esquemas dejan un residuo muy parecido (~1.550 GWh), "
           "así que probamos las dos variantes más adelante con Holt-Winters.")
p2.font.size = Pt(16)
p2.font.italic = True
p2.font.color.rgb = GRAY
add_footer(slide, 9)
set_notes(slide, """
Para entender mejor la serie, la "desarmamos" en tres piezas: una tendencia de largo
plazo, un patrón estacional que se repite cada año, y un residuo -lo que queda sin
explicar-. Hay dos formas de combinar esas piezas: sumándolas (esquema aditivo, a la
izquierda) o multiplicándolas (esquema multiplicativo, a la derecha).
La diferencia práctica es si el tamaño de los picos estacionales se mantiene constante
en GWh (aditivo) o si crece en proporción al nivel de la serie (multiplicativo). En
nuestro caso, ambos esquemas dejan residuos muy similares, por lo que no descartamos
ninguno de entrada: más adelante probamos las dos variantes con Holt-Winters y dejamos
que los resultados decidan.
""")

# ============================================================== SLIDE 10 — PATRÓN ESTACIONAL Y MEDIAS MÓVILES
slide = add_slide()
add_title_bar(slide, "El patrón que se repite cada año", "3.1 Estacionalidad y Tendencia")
add_picture_fit(slide, f"{ASSETS}/22_patron_estacional.png", Inches(0.4), Inches(1.3), Inches(6.2), Inches(4.5))
add_picture_fit(slide, f"{ASSETS}/24_medias_moviles.png", Inches(6.7), Inches(1.3), Inches(6.2), Inches(4.5))
note = slide.shapes.add_textbox(Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.2))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Izquierda: el promedio de generación por mes muestra dos picos -verano (enero/febrero, por "
          "el aire acondicionado) e invierno (junio/julio, por calefacción)- con valles en las épocas de transición.")
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p2 = tf.add_paragraph()
p2.text = ("Derecha: suavizando la serie con un promedio móvil de 12 meses, se ve un nivel estable "
           "entre 2015 y 2021, y la caída sostenida de 2022-2023.")
p2.font.size = Pt(16)
p2.font.italic = True
p2.font.color.rgb = GRAY
add_footer(slide, 10)
set_notes(slide, """
Este gráfico de la izquierda es clave para entender todo el trabajo: muestra el promedio
de generación para cada mes del año, y se ve clarísimo un patrón "bimodal", con dos
picos -uno en verano, por el uso de aire acondicionado, y otro en invierno, por la
calefacción-, con valles en otoño y primavera.
Este es justamente el patrón que un modelo sin componente estacional, como ARIMA simple,
no puede captar. A la derecha, suavizando la serie con un promedio de 12 meses, se ve el
nivel general: bastante estable hasta 2021, y con la caída sostenida que ya mencionamos
en 2022-2023.
""")

# ============================================================== SLIDE 11 — RESIDUOS DE LA DESCOMPOSICIÓN
slide = add_slide()
add_title_bar(slide, "¿Quedó algo sin explicar?", "3.1 Diagnóstico de los Residuos")
add_picture_fit(slide, f"{ASSETS}/26_residuos_descomposicion.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(4.6))
add_picture_fit(slide, f"{ASSETS}/28_acf_residuos_descomposicion.png", Inches(8.2), Inches(1.6), Inches(4.7), Inches(2.4))

tbl_shape = slide.shapes.add_table(3, 3, Inches(8.2), Inches(4.25), Inches(4.7), Inches(1.6))
table = tbl_shape.table
fill_table_text(table, [
    ["Chequeo", "Resultado", "p-valor"],
    ["¿Son normales?", "Sí", "0,184"],
    ["¿Son estables?", "Sí", "0,000"],
])
style_table(table, font_size=15, header_font_size=15, col_widths=[Inches(2.1), Inches(1.3), Inches(1.3)],
            header_height=Inches(0.55), row_height=Inches(0.5))

note = slide.shapes.add_textbox(Inches(8.2), Inches(6.05), Inches(4.7), Inches(1.2))
p = note.text_frame.paragraphs[0]
p.text = "Lo que sobra después de sacar tendencia y estacionalidad se comporta como \"ruido\": no tiene patrones escondidos. Buena señal."
p.font.size = Pt(15)
p.font.italic = True
p.font.color.rgb = GRAY
note.text_frame.word_wrap = True
add_footer(slide, 11)
set_notes(slide, """
Una vez que separamos tendencia y estacionalidad, queda un "resto" -el residuo-. La
pregunta es: ¿ese resto tiene algún patrón escondido que se nos pasó, o es simplemente
ruido sin estructura?
Hicimos dos chequeos: uno de normalidad, que confirma que los residuos se distribuyen
de forma normal; y uno de estacionariedad sobre esos residuos, que confirma que son
estables en el tiempo. En el gráfico de autocorrelación de la derecha tampoco se ven
patrones llamativos. Conclusión: la descomposición en tendencia + estacionalidad explica
bien la serie, y lo que queda es efectivamente "ruido".
""")

# ============================================================== SLIDE 12 — ESTACIONARIEDAD ADF/KPSS
slide = add_slide()
add_title_bar(slide, "¿La serie es \"estable\"? (ADF y KPSS)", "3.2 Pruebas de Estacionariedad")

items = [
    "Para usar ARIMA/SARIMA, la serie tiene que ser \"estable\" en el tiempo (estacionaria)",
    "Usamos dos pruebas que se complementan: si las dos coinciden, la conclusión es más confiable",
]
add_bullets(slide, items, Inches(0.8), Inches(1.35), Inches(11.7), Inches(1.3), font_size=18)

tbl_shape = slide.shapes.add_table(5, 5, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.6))
table = tbl_shape.table
fill_table_text(table, [
    ["Versión de la serie", "Prueba 1 (ADF)", "Prueba 2 (KPSS)", "", "¿Es estable?"],
    ["Serie original", "No pasa", "Pasa", "", "No del todo"],
    ["Con 1 diferencia simple (d=1)", "Pasa", "Pasa", "", "Sí"],
    ["Con 1 diferencia estacional (D=1)", "No pasa", "Pasa", "", "No del todo"],
    ["Con ambas diferencias (d=1 y D=1)", "Pasa", "Pasa", "", "Sí, la mejor opción"],
])
# merge col 3-4 visually by removing border between them — simpler: just style normally
style_table(table, font_size=16, header_font_size=16,
            col_widths=[Inches(3.6), Inches(2.0), Inches(2.0), Inches(0.3), Inches(3.8)],
            header_height=Inches(0.6), row_height=Inches(0.5))
highlight_row(table, 4, GOOD_BG)

note = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.6))
p = note.text_frame.paragraphs[0]
p.text = ("Conclusión: hace falta aplicar las dos diferenciaciones -una simple y una estacional- "
          "para que la serie quede \"estable\". Por eso el SARIMA que construimos usa d=1 y D=1.")
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = NAVY
note.text_frame.word_wrap = True
add_footer(slide, 12)
set_notes(slide, """
Antes de armar el ARIMA y el SARIMA, tenemos que chequear si la serie es "estable" en el
tiempo -en la jerga, "estacionaria"-. Para eso usamos dos pruebas estadísticas distintas
que se complementan: si ambas coinciden en la misma conclusión, podemos confiar más en
el resultado.
La tabla resume lo que encontramos: la serie original no pasa del todo la primera prueba.
Aplicando una diferenciación simple -comparar cada mes con el anterior- mejora bastante.
Aplicando solo la diferenciación estacional -comparar cada mes con el mismo mes del año
anterior- no alcanza. Pero combinando las dos diferenciaciones, la serie queda estable
según ambas pruebas. Por eso el SARIMA que armamos usa esas dos diferenciaciones, lo que
en la jerga técnica se anota como d=1 y D=1.
""")

# ============================================================== SLIDE 13 — MODELOS SELECCIONADOS
slide = add_slide()
add_title_bar(slide, "Los 4 modelos que vamos a comparar", "3.3 Resumen de Modelos")
tbl_shape = slide.shapes.add_table(6, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.6))
table = tbl_shape.table
fill_table_text(table, [
    ["Modelo", "¿Qué lo distingue?", "¿Tiene en cuenta la estacionalidad?"],
    ["ARIMA", "El más simple: solo mira los valores recientes. Nuestro \"punto de partida\"", "No"],
    ["SARIMA", "Como ARIMA, pero agrega el ciclo de 12 meses", "Sí"],
    ["Holt-Winters (aditivo)", "Suaviza nivel, tendencia y estacionalidad como cantidades fijas", "Sí"],
    ["Holt-Winters (multiplicativo)", "Igual, pero la estacionalidad es un % del nivel", "Sí"],
    ["Prophet", "Modelo de Meta/Facebook, ajusta tendencia y estacionalidad automáticamente", "Sí (anual)"],
])
style_table(table, font_size=16, header_font_size=17,
            col_widths=[Inches(3.0), Inches(6.3), Inches(2.4)],
            header_height=Inches(0.6), row_height=Inches(0.58))

note = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5))
p = note.text_frame.paragraphs[0]
p.text = ("Para que la comparación sea justa: entrenamos todos los modelos con los datos de "
          "2015-2022, y los evaluamos prediciendo 2023 (12 meses que ningún modelo vio antes).")
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = GRAY
note.text_frame.word_wrap = True
add_footer(slide, 13)
set_notes(slide, """
Acá tenemos los cuatro modelos que vamos a comparar. ARIMA es el más simple y nuestro
punto de partida -no tiene en cuenta la estacionalidad-. SARIMA es básicamente ARIMA más
el ciclo de 12 meses. Holt-Winters tiene dos versiones, aditiva y multiplicativa, que se
diferencian en cómo tratan el tamaño de la estacionalidad. Y Prophet es el modelo de Meta,
que ajusta automáticamente tendencia y estacionalidad sin que tengamos que elegir tantos
parámetros a mano.
Para que la comparación sea justa, a todos los entrenamos con los datos de 2015 a 2022, y
los evaluamos viendo qué tan bien predicen 2023 -que ninguno vio durante el entrenamiento-.
""")

# ============================================================== SLIDE 14 — ARIMA
slide = add_slide()
add_title_bar(slide, "Modelo 1: ARIMA — el punto de partida", "3.3 Descripción de Modelos")
add_picture_fit(slide, f"{ASSETS}/37_forecast_arima.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(4.6))
items = [
    "Mira solo el mes anterior y los errores recientes para predecir el siguiente",
    "No sabe nada de \"verano\" o \"invierno\": no tiene componente estacional",
    "Por eso esperamos que sea el que peor prediga frente a un patrón tan estacional",
    "",
    "Sirve como \"piso\": si otro modelo no le gana a este, no vale la pena usarlo",
]
add_bullets(slide, items, Inches(8.2), Inches(1.6), Inches(4.7), Inches(4.7), font_size=17)
add_footer(slide, 14)
set_notes(slide, """
Empezamos con ARIMA, el modelo más simple de los cuatro. Básicamente mira el valor del
mes anterior y los errores recientes para proyectar el siguiente mes, pero no tiene
ningún concepto de "ciclo anual": no sabe que enero se parece a otros eneros.
Por eso esperamos que sea el que peor desempeño tenga, justamente porque nuestra serie
tiene una estacionalidad tan marcada. Lo incluimos como punto de referencia: si los demás
modelos no le ganan claramente a ARIMA, no estaría justificado usar algo más complejo.
En el gráfico se puede ver que su pronóstico para 2023 es bastante "plano", sin replicar
los picos de verano e invierno.
""")

# ============================================================== SLIDE 15 — SARIMA
slide = add_slide()
add_title_bar(slide, "Modelo 2: SARIMA — ARIMA + estacionalidad", "3.3 Descripción de Modelos")
add_picture_fit(slide, f"{ASSETS}/39_forecast_sarima.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(4.6))
items = [
    "Es ARIMA, pero además compara cada mes con el mismo mes del año anterior",
    "Las dos diferenciaciones (d=1 y D=1) que vimos antes están incorporadas acá",
    "Es el modelo \"candidato natural\" para una serie tan estacional como la nuestra",
    "",
    "Esperamos que sea de los mejores, ya que ataca directamente el patrón anual",
]
add_bullets(slide, items, Inches(8.2), Inches(1.6), Inches(4.7), Inches(4.7), font_size=17)
add_footer(slide, 15)
set_notes(slide, """
SARIMA es, en esencia, ARIMA potenciado con un componente estacional. Además de mirar el
mes anterior, compara cada mes con el mismo mes del año pasado -por ejemplo, este enero
con el enero anterior-. Las dos diferenciaciones que justificamos con las pruebas ADF y
KPSS están incorporadas en este modelo.
Por construcción, este es el candidato más lógico para una serie con un patrón estacional
tan fuerte como el nuestro, y en el gráfico ya se nota: el pronóstico para 2023 sigue la
forma de los picos de verano e invierno, a diferencia del pronóstico plano de ARIMA.
""")

# ============================================================== SLIDE 16 — HOLT-WINTERS
slide = add_slide()
add_title_bar(slide, "Modelo 3: Holt-Winters — suavizado inteligente", "3.3 Descripción de Modelos")
add_picture_fit(slide, f"{ASSETS}/41_forecast_hw_aditivo.png", Inches(0.4), Inches(1.3), Inches(6.2), Inches(3.6))
add_picture_fit(slide, f"{ASSETS}/43_forecast_hw_multiplicativo.png", Inches(6.7), Inches(1.3), Inches(6.2), Inches(3.6))
items = [
    "Separa la serie en 3 \"piezas\" -nivel, tendencia y estacionalidad- y las va actualizando mes a mes a medida que llegan nuevos datos",
    "Versión aditiva (izquierda): la estacionalidad suma/resta siempre la misma cantidad de GWh",
    "Versión multiplicativa (derecha): la estacionalidad es un porcentaje del nivel actual",
]
add_bullets(slide, items, Inches(0.8), Inches(5.05), Inches(11.7), Inches(1.8), font_size=17)
add_footer(slide, 16)
set_notes(slide, """
Holt-Winters funciona distinto: en lugar de buscar fórmulas matemáticas complejas, separa
la serie en tres piezas -nivel, tendencia y estacionalidad- y las va "suavizando" y
actualizando mes a mes, dándole más o menos peso a los datos recientes.
Probamos dos versiones: la aditiva, donde el efecto estacional siempre suma o resta la
misma cantidad de GWh; y la multiplicativa, donde ese efecto es un porcentaje del nivel
actual. Como vimos en la descomposición, ambas variantes eran candidatas razonables, así
que las probamos a las dos y dejamos que la comparación final decida cuál funciona mejor.
""")

# ============================================================== SLIDE 17 — PROPHET
slide = add_slide()
add_title_bar(slide, "Modelo 4: Prophet — el de Meta/Facebook", "3.3 Descripción de Modelos")
add_picture_fit(slide, f"{ASSETS}/45_forecast_prophet.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(4.6))
items = [
    "Divide la serie en tendencia + estacionalidad, igual que los anteriores, pero ajusta todo automáticamente",
    "Usa una técnica matemática (series de Fourier) para representar la estacionalidad anual",
    "No necesitamos elegir tantos parámetros a mano ni verificar estacionariedad",
    "",
    "Lo incluimos como cuarto modelo para tener un punto de comparación \"moderno\" y de uso muy extendido en la industria",
]
add_bullets(slide, items, Inches(8.2), Inches(1.4), Inches(4.7), Inches(5.2), font_size=16)
add_footer(slide, 17)
set_notes(slide, """
El cuarto modelo es Prophet, desarrollado por Meta y muy usado en la industria para
pronósticos de negocio. También separa la serie en tendencia y estacionalidad, pero lo
hace de forma automática, usando una técnica matemática llamada series de Fourier para
representar el ciclo anual.
La ventaja es que no tenemos que elegir tantos parámetros a mano ni preocuparnos por la
estacionariedad. Lo incluimos como un cuarto modelo de referencia, muy usado en la
práctica, para ver cómo se compara contra los enfoques más "clásicos" como SARIMA y
Holt-Winters.
""")

# ============================================================== SLIDE 18 — COMPARACIÓN DE MÉTRICAS
slide = add_slide()
add_title_bar(slide, "¡El momento de la verdad! ¿Quién ganó?", "4.1 Comparación de Modelos (2023)")
add_picture_fit(slide, f"{ASSETS}/47_comparacion_modelos.png", Inches(0.4), Inches(1.3), Inches(6.8), Inches(4.6))

tbl_shape = slide.shapes.add_table(6, 2, Inches(7.4), Inches(1.5), Inches(5.5), Inches(3.2))
table = tbl_shape.table
fill_table_text(table, [
    ["Modelo", "Error promedio (MAPE)"],
    ["SARIMA", "8,15 %"],
    ["Holt-Winters (multiplicativo)", "9,02 %"],
    ["Holt-Winters (aditivo)", "9,32 %"],
    ["Prophet", "11,61 %"],
    ["ARIMA", "13,01 %"],
])
style_table(table, font_size=16, header_font_size=17, col_widths=[Inches(3.6), Inches(1.9)],
            header_height=Inches(0.6), row_height=Inches(0.5))
highlight_row(table, 1, GOOD_BG)

note = slide.shapes.add_textbox(Inches(7.4), Inches(4.85), Inches(5.5), Inches(2.2))
p = note.text_frame.paragraphs[0]
p.text = ("El MAPE es el error promedio en porcentaje: \"en promedio, el modelo se equivoca "
          "un X% del valor real\". Cuanto más bajo, mejor.")
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p2 = note.text_frame.add_paragraph()
p2.text = ("Ganó SARIMA con 8,15% de error. La diferencia con ARIMA (13,01%) muestra "
           "cuánto aporta modelar la estacionalidad.")
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = NAVY
note.text_frame.word_wrap = True
add_footer(slide, 18)
set_notes(slide, """
Este es el resultado clave de todo el trabajo. Para comparar los modelos usamos el MAPE,
que es el error promedio expresado en porcentaje: si el MAPE es 8%, en promedio el
modelo se equivoca un 8% del valor real. Es una métrica fácil de comunicar porque no
depende de la unidad -GWh-, sino que es relativa.
El ganador fue SARIMA, con un error de 8,15%. Le siguen las dos versiones de
Holt-Winters, alrededor del 9%, después Prophet con 11,6%, y por último ARIMA con 13%.
Este resultado confirma nuestra hipótesis inicial: el modelo que mejor capta el patrón
estacional -SARIMA- es el que predice mejor, y la diferencia con ARIMA -que no lo tiene
en cuenta- es de casi 5 puntos porcentuales. Usamos también otras dos métricas, MAE y
RMSE, y el orden de los modelos es el mismo en las tres, lo que le da más solidez a esta
conclusión.
""")

# ============================================================== SLIDE 19 — SARIMA VS REAL
slide = add_slide()
add_title_bar(slide, "Mirando de cerca al ganador: SARIMA", "4.2 SARIMA vs. la Realidad")
add_picture_fit(slide, f"{ASSETS}/51_sarima_vs_real.png", Inches(0.4), Inches(1.3), Inches(6.2), Inches(4.7))
add_picture_fit(slide, f"{ASSETS}/58_real_vs_predicho_sarima.png", Inches(6.7), Inches(1.3), Inches(6.2), Inches(4.7))
note = slide.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(12.1), Inches(1.0))
p = note.text_frame.paragraphs[0]
p.text = ("La línea de predicción de SARIMA sigue de cerca a la línea real, incluyendo los picos "
          "estacionales. Las mayores diferencias aparecen en los meses donde la caída de 2023 fue más fuerte.")
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = GRAY
note.text_frame.word_wrap = True
add_footer(slide, 19)
set_notes(slide, """
Acá podemos ver en detalle qué tan bien predijo SARIMA comparado con lo que realmente
pasó en 2023. La línea de predicción sigue de cerca a la línea real, incluyendo la forma
de los picos de verano e invierno -algo que, como vimos, ARIMA no lograba-.
Las mayores diferencias entre predicción y realidad aparecen en los meses donde la
caída de generación de 2023 fue más pronunciada, justamente porque el modelo se entrenó
con datos hasta 2022 y no podía "saber" que venía esa baja tan marcada.
""")

# ============================================================== SLIDE 20 — DIAGNÓSTICO RESIDUOS SARIMA
slide = add_slide()
add_title_bar(slide, "¿SARIMA está bien armado?", "4.3 Validación de Residuos")
add_picture_fit(slide, f"{ASSETS}/53_acf_pacf_residuos_sarima.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(3.0))

tbl_shape = slide.shapes.add_table(5, 3, Inches(8.2), Inches(1.5), Inches(4.7), Inches(2.4))
table = tbl_shape.table
fill_table_text(table, [
    ["Horizonte", "¿Hay patrones sin explicar?", "p-valor"],
    ["6 meses", "No", "0,946"],
    ["12 meses", "No", "0,125"],
    ["18 meses", "No", "0,302"],
    ["24 meses", "No", "0,639"],
])
style_table(table, font_size=15, header_font_size=15, col_widths=[Inches(1.6), Inches(2.1), Inches(1.0)],
            header_height=Inches(0.6), row_height=Inches(0.45))

note = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(12.4), Inches(2.0))
p = note.text_frame.paragraphs[0]
p.text = ("Esta prueba (Ljung-Box) chequea si quedó algún patrón sin explicar en lo que el modelo "
          "no pudo predecir. En todos los horizontes -incluso a 12 y 24 meses, los \"ciclos\" anuales- "
          "la respuesta es \"no\": SARIMA capturó bien tanto el corto plazo como el ciclo anual.")
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = NAVY
note.text_frame.word_wrap = True
add_footer(slide, 20)
set_notes(slide, """
Que SARIMA haya tenido el menor error no es suficiente: también queremos confirmar que
está "bien armado", es decir, que no dejó información sin aprovechar. Para eso miramos
los residuos -la diferencia entre lo predicho y lo real- y les hacemos una prueba
llamada Ljung-Box, que chequea si queda algún patrón escondido.
La tabla muestra que, para todos los horizontes que probamos -incluso a 12 y 24 meses,
que son justamente los "ciclos" anuales-, la respuesta es que no quedan patrones sin
explicar. Esto confirma que SARIMA no solo predijo mejor, sino que está correctamente
especificado: no hay una versión "más simple que se nos escapó" que pudiera mejorar el
resultado.
""")

# ============================================================== SLIDE 21 — AIC GRID SEARCH
slide = add_slide()
add_title_bar(slide, "¿Y si lo elegía la computadora?", "4.4 Extensión 1: Selección Automática")

tbl_shape = slide.shapes.add_table(6, 2, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.3))
table = tbl_shape.table
fill_table_text(table, [
    ["Configuración SARIMA", "Puntaje (AIC, menor = mejor ajuste)"],
    ["(1,1,1)(0,1,1,12)  ← elegido por la PC", "1244,08"],
    ["(1,1,1)(1,1,1,12)  ← el nuestro", "1245,97"],
    ["(0,1,1)(0,1,1,12)", "1248,39"],
    ["(0,1,1)(1,1,1,12)", "1250,38"],
    ["(1,0,1)(0,1,1,12)", "1261,68"],
])
style_table(table, font_size=14, header_font_size=14, col_widths=[Inches(3.6), Inches(2.2)],
            header_height=Inches(0.7), row_height=Inches(0.5))

tbl_shape2 = slide.shapes.add_table(3, 2, Inches(7.0), Inches(1.5), Inches(5.9), Inches(1.8))
table2 = tbl_shape2.table
fill_table_text(table2, [
    ["Configuración", "Error (MAPE) prediciendo 2023"],
    ["Elegida por la PC", "9,47 %"],
    ["La nuestra (manual)", "8,15 %"],
])
style_table(table2, font_size=15, header_font_size=15, col_widths=[Inches(3.3), Inches(2.6)],
            header_height=Inches(0.65), row_height=Inches(0.6))
highlight_row(table2, 2, GOOD_BG)

note = slide.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.9), Inches(2.0))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Probamos también una búsqueda automática que recorre muchas combinaciones y elige la de mejor \"puntaje\" interno (AIC)."
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = ("La diferencia de puntaje es mínima, pero a la hora de predecir 2023 de verdad, "
           "nuestra configuración elegida \"a mano\" predijo mejor (8,15% vs. 9,47%). Moraleja: "
           "un buen puntaje interno no garantiza el mejor pronóstico real, por eso siempre hay "
           "que validar contra datos reales.")
p2.font.size = Pt(17)
p2.font.color.rgb = GRAY
p2.space_before = Pt(8)
add_footer(slide, 21)
set_notes(slide, """
Como extensión, probamos qué pasaría si dejábamos que una búsqueda automática eligiera
la configuración de SARIMA, en lugar de elegirla nosotros. Esa búsqueda recorre muchas
combinaciones posibles y elige la que tiene mejor "puntaje" interno, llamado AIC, que
mide qué tan bien se ajusta el modelo a los datos de entrenamiento penalizando la
complejidad.
La configuración elegida automáticamente tuvo un puntaje apenas mejor que la nuestra.
Pero cuando las hicimos competir prediciendo 2023 -los datos reales que nos interesan-,
nuestra configuración elegida a mano predijo mejor: 8,15% de error contra 9,47%. La
moraleja es importante: un buen puntaje "interno" no garantiza el mejor pronóstico en la
práctica, por eso es clave validar siempre contra datos reales que el modelo no vio.
""")

# ============================================================== SLIDE 22 — GARCH
slide = add_slide()
add_title_bar(slide, "Bonus: ¿qué tan \"nerviosa\" es la serie?", "4.5 Extensión 2: Volatilidad (GARCH)")
add_picture_fit(slide, f"{ASSETS}/63_garch_volatilidad.png", Inches(0.4), Inches(1.3), Inches(7.2), Inches(4.7))

tbl_shape = slide.shapes.add_table(4, 3, Inches(7.9), Inches(1.5), Inches(5.0), Inches(2.0))
table = tbl_shape.table
fill_table_text(table, [
    ["Parámetro", "¿Qué mide?", "¿Es relevante?"],
    ["ω", "Nivel \"normal\" de volatilidad", "Sí"],
    ["α (shocks)", "Reacción a una sorpresa puntual", "No tanto"],
    ["β (memoria)", "Cuánto dura un período movido", "Sí"],
])
style_table(table, font_size=14, header_font_size=14,
            col_widths=[Inches(1.6), Inches(2.4), Inches(1.0)],
            header_height=Inches(0.6), row_height=Inches(0.5))

note = slide.shapes.add_textbox(Inches(7.9), Inches(3.85), Inches(5.0), Inches(3.2))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Todos los modelos anteriores predicen el valor promedio. GARCH va un paso más allá: "
          "predice qué tan \"ancho\" debería ser el margen de error mes a mes.")
p.font.size = Pt(17)
p.font.color.rgb = GRAY
p2 = tf.add_paragraph()
p2.text = ("Encontramos que los períodos movidos tienden a continuar un tiempo (hay \"memoria\"), "
           "pero un shock puntual no dispara por sí solo un período de alta volatilidad.")
p2.font.size = Pt(17)
p2.font.color.rgb = GRAY
p2.space_before = Pt(10)
p3 = tf.add_paragraph()
p3.text = "En la práctica: permite dar intervalos de confianza más realistas, más anchos cuando \"está movido\" y más angostos cuando está tranquilo."
p3.font.size = Pt(17)
p3.font.bold = True
p3.font.color.rgb = NAVY
p3.space_before = Pt(10)
add_footer(slide, 22)
set_notes(slide, """
Como bonus, agregamos un análisis de volatilidad con un modelo llamado GARCH. La idea es
distinta a todo lo anterior: en lugar de predecir el valor promedio, predice qué tan
"ancho" debería ser el margen de error en cada momento -es decir, mide qué tan "nerviosa"
está la serie-.
Encontramos que hay cierta "memoria" en la volatilidad: si un mes es movido, es más
probable que el siguiente también lo sea. Pero un shock puntual aislado no necesariamente
dispara por sí solo un período de alta volatilidad. En la práctica, esto nos permite dar
intervalos de confianza más realistas: más anchos en los períodos movidos, y más angostos
en los tranquilos, en lugar de un margen de error fijo para todos los meses.
""")

# ============================================================== SLIDE 23 — PRONÓSTICO 2024
slide = add_slide()
add_title_bar(slide, "¿Y qué esperamos para 2024?", "4.6 Pronóstico Operativo")
add_picture_fit(slide, f"{ASSETS}/65_pronostico_2024.png", Inches(0.4), Inches(1.3), Inches(7.6), Inches(4.7))

tbl_shape = slide.shapes.add_table(7, 4, Inches(8.2), Inches(1.3), Inches(4.7), Inches(4.7))
table = tbl_shape.table
fill_table_text(table, [
    ["Mes", "Pronóstico", "Mínimo esperado", "Máximo esperado"],
    ["Enero", "35.332", "31.161", "39.503"],
    ["Febrero", "29.570", "24.767", "34.373"],
    ["Marzo", "33.997", "28.937", "39.057"],
    ["Junio", "32.933", "27.512", "38.353"],
    ["Julio", "32.693", "27.182", "38.203"],
    ["Diciembre", "30.696", "24.772", "36.619"],
])
style_table(table, font_size=13, header_font_size=14,
            col_widths=[Inches(1.1), Inches(1.2), Inches(1.2), Inches(1.2)],
            header_height=Inches(0.7), row_height=Inches(0.55))

note = slide.shapes.add_textbox(Inches(8.2), Inches(6.25), Inches(4.7), Inches(1.0))
p = note.text_frame.paragraphs[0]
p.text = "Valores en GWh. Tabla resumida (6 de 12 meses) — el detalle completo está en el informe."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY
note.text_frame.word_wrap = True
add_footer(slide, 23)
set_notes(slide, """
Por último, usamos SARIMA -nuestro modelo ganador- reentrenado con todos los datos
disponibles, de 2015 a 2023, para proyectar los 12 meses de 2024. El gráfico muestra esa
proyección junto con una banda de confianza del 95%.
Como era de esperar, el pronóstico reproduce el mismo patrón bimodal que vimos en toda la
serie: picos en enero y en junio-julio. La tabla muestra, para algunos meses, el valor
esperado y el rango -mínimo y máximo- dentro del cual esperamos que caiga el valor real
con un 95% de confianza. Esta es la salida que un equipo de planificación energética
podría usar como insumo para sus decisiones.
""")

# ============================================================== SLIDE 24 — CONCLUSIONES 1
slide = add_slide()
add_title_bar(slide, "Lo que aprendimos (parte 1)", "5. Conclusiones")
items = [
    "Respondiendo la pregunta inicial: SARIMA fue el mejor (8,15% de error), y la estacionalidad anual resultó ser el factor más importante para predecir bien.",
    "El orden final fue: SARIMA, después Holt-Winters (~9%), después Prophet (11,6%) y por último ARIMA (13%) — confirmando que cuanto mejor se modela la estacionalidad, mejor se predice.",
    "Entre las dos versiones de Holt-Winters, la multiplicativa fue apenas mejor (9,02% vs. 9,32%) — una diferencia pequeña, esperable según lo que ya habíamos visto en la descomposición.",
    "La fuerte caída de 2022-2023 fue el principal \"dolor de cabeza\" para todos los modelos: ninguno podía anticipar ese quiebre sin información externa (ej. actividad económica).",
]
add_bullets(slide, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8), font_size=19)
add_footer(slide, 24)
set_notes(slide, """
Para cerrar, repasemos las conclusiones principales. La pregunta que nos hicimos al
principio queda respondida: SARIMA fue el modelo más preciso, con un error de 8,15%, y
esto confirma que la estacionalidad anual es el factor más determinante para predecir
bien esta serie.
El orden final -SARIMA, Holt-Winters, Prophet y ARIMA- tiene una lectura clara: cuanto
mejor modela un método la estacionalidad, mejor predice. Entre las dos versiones de
Holt-Winters la diferencia fue mínima, algo que ya habíamos anticipado en la
descomposición. Y un punto importante: la fuerte caída de generación en 2022-2023 fue el
principal desafío para todos los modelos, porque ninguno podía "adivinar" ese cambio sin
información externa, como por ejemplo el nivel de actividad económica del país.
""")

# ============================================================== SLIDE 25 — CONCLUSIONES 2
slide = add_slide()
add_title_bar(slide, "Lo que aprendimos (parte 2)", "5. Conclusiones y Aprendizajes")
items = [
    "El \"piloto automático\" (selección por AIC) no siempre gana: nuestra configuración elegida a mano predijo mejor en datos reales (8,15% vs. 9,47%).",
    "GARCH mostró que la volatilidad \"tiene memoria\": los períodos movidos tienden a continuar, lo que permite dar márgenes de error más realistas.",
    "Validamos que SARIMA no solo ganó en precisión, sino que está bien especificado: no deja patrones sin explicar en sus errores.",
    "Aprendizaje general: ningún indicador por sí solo (ni el puntaje interno, ni el error de predicción, ni las pruebas de residuos) alcanza — hay que mirar los tres juntos para elegir con confianza.",
]
add_bullets(slide, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8), font_size=19)
add_footer(slide, 25)
set_notes(slide, """
Algunas reflexiones finales. Primero, vimos que el "piloto automático" de selección de
modelos no siempre acierta: la configuración elegida a mano predijo mejor en datos
reales que la elegida automáticamente por puntaje interno.
Segundo, el análisis de volatilidad con GARCH mostró que los períodos "movidos" tienden a
durar un tiempo, lo que en la práctica permite dar márgenes de error más realistas y no
un número fijo para todos los meses.
Tercero, no solo nos quedamos con que SARIMA tuvo el menor error: también confirmamos que
está bien especificado, sin patrones sin explicar en sus errores.
Y como aprendizaje general del trabajo: ningún indicador por sí solo es suficiente. Hay
que combinar el ajuste interno del modelo, el error prediciendo datos reales, y el
diagnóstico de los residuos para elegir un modelo con confianza. Con esto cerramos la
presentación, ¡gracias!
""")

# ============================================================== SLIDE 26 — GRACIAS
slide = add_slide()
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = PAGE_BG
bg.line.fill.background()

# formas decorativas, mismo lenguaje visual que la portada
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(-1.6), Inches(4.0), Inches(4.0))
circle1.fill.solid()
circle1.fill.fore_color.rgb = NAVY
circle1.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, SW - Inches(2.6), SH - Inches(2.6), Inches(4.0), Inches(4.0))
circle2.fill.solid()
circle2.fill.fore_color.rgb = TEAL
circle2.line.fill.background()
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.55), SW, Pt(2))
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()

if os.path.exists(LOGO_PATH):
    plate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, (SW - Inches(2.3)) / 2, Inches(0.55), Inches(2.3), Inches(0.95))
    plate.fill.solid()
    plate.fill.fore_color.rgb = WHITE
    plate.line.color.rgb = BORDER
    plate.line.width = Pt(1)
    plate.shadow.inherit = False
    add_picture_fit(slide, LOGO_PATH, (SW - Inches(2.0)) / 2, Inches(0.65), Inches(2.0), Inches(0.75))

tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.65), Inches(11.5), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Muchas gracias"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Preguntas y discusión"
p2.font.size = Pt(24)
p2.font.color.rgb = TEAL
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)
p3 = tf.add_paragraph()
p3.text = "Gustavo Rivas · Carlos Rivas · Fermín Rodríguez"
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)
p4 = tf.add_paragraph()
p4.text = "Análisis de Series de Tiempo (02MIA2026) — Especialización en IA, UBA FIUBA"
p4.font.size = Pt(13)
p4.font.color.rgb = GRAY
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(6)
set_notes(slide, """
Muchas gracias por la atención. Quedamos abiertos a preguntas sobre cualquiera de los
modelos, los resultados o las extensiones que presentamos.
""")

prs.save("INFORME_PRESENTACION.pptx")
print("OK -", len(prs.slides._sldIdLst), "slides")
